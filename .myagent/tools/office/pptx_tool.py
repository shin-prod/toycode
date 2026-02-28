"""PowerPoint (PPTX) 操作ツール。python-pptx を使用する。"""

import json
import os
from typing import Optional

from utils.file_guard import check_write_allowed, resolve_path
from utils.logger import get_logger

# shape_type 文字列 → MSO_AUTO_SHAPE_TYPE 整数値マッピング
_SHAPE_MAP: dict[str, int] = {
    "rectangle": 1,
    "rounded_rectangle": 5,
    "oval": 9,
    "ellipse": 9,
    "triangle": 7,
    "diamond": 4,
    "hexagon": 10,
    "pentagon": 56,
    "right_arrow": 13,
    "left_arrow": 34,
    "up_arrow": 35,
    "down_arrow": 36,
    "double_arrow": 37,
    "star4": 11,
    "star5": 12,
    "callout": 57,
}

_ALIGN_MAP = {
    "center": None,   # PP_ALIGN.CENTER
    "left": None,     # PP_ALIGN.LEFT
    "right": None,    # PP_ALIGN.RIGHT
    "justify": None,  # PP_ALIGN.JUSTIFY
}

logger = get_logger(__name__)


def _get_align(text_align: str):
    """PP_ALIGN 列挙値を返す（遅延インポート）。"""
    from pptx.enum.text import PP_ALIGN  # noqa: PLC0415
    return {
        "center": PP_ALIGN.CENTER,
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(text_align.lower(), PP_ALIGN.LEFT)


def _parse_color(hex_color: str):
    """16進数カラー文字列を RGBColor に変換する。"""
    from pptx.dml.color import RGBColor  # noqa: PLC0415
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _apply_text_frame(
    tf,
    text: str,
    font_size: int,
    bold: bool,
    font_color: str,
    text_align: str,
    auto_size_mode: str = "fit_text",
) -> None:
    """テキストフレームにテキスト・書式を適用する共通ヘルパー。

    Args:
        tf: python-pptx の TextFrame オブジェクト
        text: 設定するテキスト（\\n で段落分割）
        font_size: フォントサイズ (pt)
        bold: 太字フラグ
        font_color: 16進数 RGB カラー文字列（空文字で無指定）
        text_align: "left" / "center" / "right" / "justify"
        auto_size_mode:
            "fit_text"  → TEXT_TO_FIT_SHAPE（テキストを縮小して図形内に収める）
            "fit_shape" → SHAPE_TO_FIT_TEXT（図形をテキストに合わせて拡大）
            "none"      → 自動サイズ調整なし
    """
    from pptx.util import Pt  # noqa: PLC0415
    from pptx.enum.text import MSO_AUTO_SIZE  # noqa: PLC0415

    _AUTO_SIZE = {
        "fit_text": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
        "fit_shape": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        "none": MSO_AUTO_SIZE.NONE,
    }

    tf.word_wrap = True
    tf.auto_size = _AUTO_SIZE.get(auto_size_mode, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    lines = text.split("\n")
    align = _get_align(text_align)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if font_color:
            run.font.color.rgb = _parse_color(font_color)


def pptx_read(path: str) -> str:
    """PPTX ファイルのスライド内容を JSON 形式で取得する。

    スライドサイズ（インチ）も返すので、要素配置の参考にすること。

    Args:
        path: PPTX ファイルのパス

    Returns:
        スライド構造の JSON 文字列
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return "エラー: python-pptx がインストールされていません。pip install python-pptx"

    resolved = resolve_path(path)
    if not os.path.isfile(resolved):
        return f"エラー: ファイルが存在しません: {resolved}"

    logger.debug("pptx_read: %s", resolved)
    try:
        prs = Presentation(resolved)
        slide_w = round(prs.slide_width / 914400, 3)   # EMU → inches
        slide_h = round(prs.slide_height / 914400, 3)

        slides_data = []
        for i, slide in enumerate(prs.slides):
            slide_info: dict = {"index": i, "shapes": []}
            for shape in slide.shapes:
                shape_info: dict = {
                    "shape_id": shape.shape_id,
                    "name": shape.name,
                    "left": round(shape.left / 914400, 3),
                    "top": round(shape.top / 914400, 3),
                    "width": round(shape.width / 914400, 3),
                    "height": round(shape.height / 914400, 3),
                }
                if shape.has_text_frame:
                    shape_info["text"] = shape.text_frame.text
                    shape_info["paragraphs"] = [
                        p.text for p in shape.text_frame.paragraphs
                    ]
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    shape_info["type"] = "picture"
                else:
                    shape_info["type"] = "text_or_other"
                slide_info["shapes"].append(shape_info)
            if slide.shapes.title:
                slide_info["title"] = slide.shapes.title.text
            slides_data.append(slide_info)

        return json.dumps(
            {
                "slide_width_inches": slide_w,
                "slide_height_inches": slide_h,
                "slide_count": len(slides_data),
                "slides": slides_data,
            },
            ensure_ascii=False,
            indent=2,
        )
    except Exception as e:
        logger.error("pptx_read エラー: %s", e)
        return f"エラー: PPTX 読み込み失敗: {e}"


def pptx_edit_slide(
    path: str,
    slide_index: int,
    title: Optional[str] = None,
    content: Optional[str] = None,
) -> str:
    """スライドのタイトル・本文テキストを変更して保存する。

    Args:
        path: PPTX ファイルのパス
        slide_index: 編集するスライドのインデックス（0始まり）
        title: 新しいタイトル（省略可）
        content: 新しい本文テキスト（\\n で箇条書き行を区切り可）

    Returns:
        成功メッセージまたはエラーメッセージ
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "エラー: python-pptx がインストールされていません。pip install python-pptx"

    resolved = resolve_path(path)
    if err := check_write_allowed(resolved):
        return err
    if not os.path.isfile(resolved):
        return f"エラー: ファイルが存在しません: {resolved}"

    logger.debug("pptx_edit_slide: %s slide=%d", resolved, slide_index)
    try:
        prs = Presentation(resolved)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return (
                f"エラー: スライドインデックスが範囲外です "
                f"(0〜{len(prs.slides) - 1})"
            )
        slide = prs.slides[slide_index]

        if title is not None:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                title_shape.text_frame.text = title
            else:
                return "エラー: このスライドにはタイトルシェイプがありません"

        if content is not None:
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    tf = shape.text_frame
                    tf.clear()
                    lines = content.split("\n")
                    for j, line in enumerate(lines):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.text = line
                    break
            else:
                return "エラー: 本文を書き込めるシェイプが見つかりませんでした"

        prs.save(resolved)
        return f"スライド {slide_index} を更新しました: {resolved}"
    except Exception as e:
        logger.error("pptx_edit_slide エラー: %s", e)
        return f"エラー: PPTX 編集失敗: {e}"


def pptx_add_slide(
    path: str,
    title: str,
    content: str = "",
) -> str:
    """PPTX ファイルにスライドを追加する。

    Args:
        path: PPTX ファイルのパス
        title: 新しいスライドのタイトル
        content: 新しいスライドの本文テキスト（\\n で複数行・箇条書き）

    Returns:
        成功メッセージまたはエラーメッセージ
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "エラー: python-pptx がインストールされていません。pip install python-pptx"

    resolved = resolve_path(path)
    if err := check_write_allowed(resolved):
        return err

    logger.debug("pptx_add_slide: %s title=%s", resolved, title)
    try:
        if os.path.isfile(resolved):
            prs = Presentation(resolved)
        else:
            os.makedirs(os.path.dirname(resolved) or ".", exist_ok=True)
            prs = Presentation()

        layout_index = 1 if len(prs.slide_layouts) > 1 else 0
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        if content:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    tf = ph.text_frame
                    tf.clear()
                    lines = content.split("\n")
                    for j, line in enumerate(lines):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.text = line
                    break

        prs.save(resolved)
        new_index = len(prs.slides) - 1
        return f"スライドを追加しました (index={new_index}): {resolved}"
    except Exception as e:
        logger.error("pptx_add_slide エラー: %s", e)
        return f"エラー: PPTX スライド追加失敗: {e}"


def pptx_add_shape(
    path: str,
    slide_index: int,
    shape_type: str,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    fill_color: str = "",
    line_color: str = "",
    font_size: int = 14,
    bold: bool = False,
    text_align: str = "center",
    font_color: str = "",
) -> str:
    """スライドにオートシェイプ（図形）を追加する。

    座標・サイズはすべてインチ単位。デフォルトスライドサイズは 10×7.5 インチ。
    テキストは図形内に自動縮小して収められる（はみ出しなし）。

    Args:
        path: PPTX ファイルのパス
        slide_index: 対象スライドのインデックス（0始まり）
        shape_type: 図形の種類。rectangle / rounded_rectangle / oval / triangle /
                    diamond / hexagon / pentagon / right_arrow / left_arrow /
                    up_arrow / down_arrow / double_arrow / star4 / star5 / callout
        left: 左端位置（インチ）
        top: 上端位置（インチ）
        width: 幅（インチ）
        height: 高さ（インチ）
        text: 図形内に表示するテキスト（\\n で改行可）（省略可）
        fill_color: 塗りつぶし色（16進数 RGB 例: "FF0000"）（省略可）
        line_color: 枠線色（16進数 RGB 例: "0000FF"）（省略可）
        font_size: テキストのフォントサイズ (pt)（デフォルト: 14）
        bold: テキストを太字にするか（デフォルト: False）
        text_align: テキスト揃え "left" / "center" / "right"（デフォルト: "center"）
        font_color: テキスト色（16進数 RGB 例: "FFFFFF"）（省略可）

    Returns:
        成功メッセージまたはエラーメッセージ
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
    except ImportError:
        return "エラー: python-pptx がインストールされていません。pip install python-pptx"

    shape_id = _SHAPE_MAP.get(shape_type.lower())
    if shape_id is None:
        return (
            f"エラー: 未知の shape_type '{shape_type}'。"
            f" 使用可能: {', '.join(_SHAPE_MAP.keys())}"
        )

    resolved = resolve_path(path)
    if err := check_write_allowed(resolved):
        return err
    if not os.path.isfile(resolved):
        return f"エラー: ファイルが存在しません: {resolved}"

    logger.debug("pptx_add_shape: %s slide=%d shape=%s", resolved, slide_index, shape_type)
    try:
        prs = Presentation(resolved)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return f"エラー: スライドインデックスが範囲外です (0〜{len(prs.slides) - 1})"
        slide = prs.slides[slide_index]

        # 座標がスライドをはみ出していないか警告
        slide_w = prs.slide_width / 914400
        slide_h = prs.slide_height / 914400
        if left + width > slide_w or top + height > slide_h:
            logger.warning(
                "図形がスライド境界を超えています: slide=(%g×%g) shape=(l=%g t=%g w=%g h=%g)",
                slide_w, slide_h, left, top, width, height,
            )

        shape = slide.shapes.add_shape(
            shape_id,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _parse_color(fill_color)

        if line_color:
            shape.line.color.rgb = _parse_color(line_color)

        if text and shape.has_text_frame:
            _apply_text_frame(
                shape.text_frame, text, font_size, bold, font_color, text_align,
                auto_size_mode="fit_text",
            )

        prs.save(resolved)
        return f"図形を追加しました (slide={slide_index}, type={shape_type}): {resolved}"
    except Exception as e:
        logger.error("pptx_add_shape エラー: %s", e)
        return f"エラー: 図形追加失敗: {e}"


def pptx_add_textbox(
    path: str,
    slide_index: int,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 16,
    bold: bool = False,
    font_color: str = "",
    text_align: str = "left",
    auto_size: bool = False,
) -> str:
    """スライドにテキストボックスを追加する。

    座標・サイズはすべてインチ単位。デフォルトスライドサイズは 10×7.5 インチ。
    テキストははみ出さないよう自動縮小される（auto_size=False のとき）。

    Args:
        path: PPTX ファイルのパス
        slide_index: 対象スライドのインデックス（0始まり）
        text: テキスト内容（\\n で改行・段落分割）
        left: 左端位置（インチ）
        top: 上端位置（インチ）
        width: 幅（インチ）
        height: 高さ（インチ）
        font_size: フォントサイズ（pt）（デフォルト: 16）
        bold: 太字にするか（デフォルト: False）
        font_color: 文字色（16進数 RGB 例: "FF0000"）（省略可）
        text_align: テキスト揃え "left" / "center" / "right"（デフォルト: "left"）
        auto_size: True のとき高さをテキストに合わせて自動拡張する（デフォルト: False）

    Returns:
        成功メッセージまたはエラーメッセージ
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return "エラー: python-pptx がインストールされていません。pip install python-pptx"

    resolved = resolve_path(path)
    if err := check_write_allowed(resolved):
        return err
    if not os.path.isfile(resolved):
        return f"エラー: ファイルが存在しません: {resolved}"

    logger.debug("pptx_add_textbox: %s slide=%d", resolved, slide_index)
    try:
        prs = Presentation(resolved)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return f"エラー: スライドインデックスが範囲外です (0〜{len(prs.slides) - 1})"
        slide = prs.slides[slide_index]

        # 座標がスライドをはみ出していないか警告
        slide_w = prs.slide_width / 914400
        slide_h = prs.slide_height / 914400
        if left + width > slide_w or top + height > slide_h:
            logger.warning(
                "テキストボックスがスライド境界を超えています: "
                "slide=(%g×%g) box=(l=%g t=%g w=%g h=%g)",
                slide_w, slide_h, left, top, width, height,
            )

        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        _apply_text_frame(
            txBox.text_frame, text, font_size, bold, font_color, text_align,
            auto_size_mode="fit_shape" if auto_size else "fit_text",
        )

        prs.save(resolved)
        return f"テキストボックスを追加しました (slide={slide_index}): {resolved}"
    except Exception as e:
        logger.error("pptx_add_textbox エラー: %s", e)
        return f"エラー: テキストボックス追加失敗: {e}"


def pptx_add_picture(
    path: str,
    slide_index: int,
    image_path: str,
    left: float,
    top: float,
    width: float = 0.0,
    height: float = 0.0,
) -> str:
    """スライドに画像を挿入する。

    Args:
        path: PPTX ファイルのパス
        slide_index: 対象スライドのインデックス（0始まり）
        image_path: 挿入する画像ファイルのパス（WORKSPACE_DIR 基準の相対パスも可）
        left: 左端位置（インチ）
        top: 上端位置（インチ）
        width: 幅（インチ）（0 の場合は height からアスペクト比維持）
        height: 高さ（インチ）（0 の場合は width からアスペクト比維持）

    Returns:
        成功メッセージまたはエラーメッセージ
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return "エラー: python-pptx がインストールされていません。pip install python-pptx"

    resolved = resolve_path(path)
    if err := check_write_allowed(resolved):
        return err
    if not os.path.isfile(resolved):
        return f"エラー: ファイルが存在しません: {resolved}"

    resolved_image = resolve_path(image_path)
    if not os.path.isfile(resolved_image):
        return f"エラー: 画像ファイルが存在しません: {resolved_image}"

    logger.debug("pptx_add_picture: %s slide=%d img=%s", resolved, slide_index, resolved_image)
    try:
        prs = Presentation(resolved)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return f"エラー: スライドインデックスが範囲外です (0〜{len(prs.slides) - 1})"
        slide = prs.slides[slide_index]

        pic_width = Inches(width) if width else None
        pic_height = Inches(height) if height else None
        slide.shapes.add_picture(
            resolved_image,
            Inches(left), Inches(top),
            width=pic_width,
            height=pic_height,
        )

        prs.save(resolved)
        return f"画像を挿入しました (slide={slide_index}): {resolved}"
    except Exception as e:
        logger.error("pptx_add_picture エラー: %s", e)
        return f"エラー: 画像挿入失敗: {e}"


def pptx_add_table(
    path: str,
    slide_index: int,
    data: list,
    left: float,
    top: float,
    width: float,
    height: float,
) -> str:
    """スライドにテーブル（表）を追加する。

    Args:
        path: PPTX ファイルのパス
        slide_index: 対象スライドのインデックス（0始まり）
        data: セルデータの2次元リスト（行×列）。1行目がヘッダーとして扱われる
        left: 左端位置（インチ）
        top: 上端位置（インチ）
        width: 幅（インチ）
        height: 高さ（インチ）

    Returns:
        成功メッセージまたはエラーメッセージ
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return "エラー: python-pptx がインストールされていません。pip install python-pptx"

    if not data or not data[0]:
        return "エラー: data が空です"

    rows = len(data)
    cols = max(len(row) for row in data)

    resolved = resolve_path(path)
    if err := check_write_allowed(resolved):
        return err
    if not os.path.isfile(resolved):
        return f"エラー: ファイルが存在しません: {resolved}"

    logger.debug(
        "pptx_add_table: %s slide=%d rows=%d cols=%d", resolved, slide_index, rows, cols
    )
    try:
        prs = Presentation(resolved)
        if slide_index < 0 or slide_index >= len(prs.slides):
            return f"エラー: スライドインデックスが範囲外です (0〜{len(prs.slides) - 1})"
        slide = prs.slides[slide_index]

        table = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top), Inches(width), Inches(height),
        ).table

        for r_idx, row in enumerate(data):
            for c_idx, cell_value in enumerate(row):
                if c_idx < cols:
                    table.cell(r_idx, c_idx).text = str(cell_value)

        prs.save(resolved)
        return f"テーブルを追加しました ({rows}行×{cols}列, slide={slide_index}): {resolved}"
    except Exception as e:
        logger.error("pptx_add_table エラー: %s", e)
        return f"エラー: テーブル追加失敗: {e}"
